from pdfminer.high_level import extract_text #pip install pdfminer.six
#pip install 'pdfminer.six[image]' to extract images
import docx2txt #pip install docx2txt
import pandas as pd
from pptx import Presentation #install python-pptx
from bs4 import BeautifulSoup
import glob
import docx #python-docx
import numpy as np
from PIL import Image #pip install pillow
#import pytesseract #pip install pytesseract
import logging
import json

class readers:
    @staticmethod
    def read_pdf(file_path):
        text=""
        text = extract_text(file_path)
        return text

    @staticmethod
    def read_txt(file_path):
        text = ""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text=file.read()
        return text

    @staticmethod
    def read_docx(file_path):
        my_text=""
        my_text = docx2txt.process(file_path, img_dir=None)
        return my_text
        
    '''
       images = glob.glob("."+"/image*.*")

     print(images)
        pytesseract.pytesseract.tesseract_cmd = r'C:\\Program Files\\Tesseract-OCR'
        for image in images:
            try:
                print(image)
                print(pytesseract.pytesseract.image_to_string(Image.open(image)))
            except:
                print("Exception")
    '''
        

    @staticmethod
    def read_xls(file_path):
        text = " "
        xls = pd.ExcelFile(file_path)
        for sheet in xls.sheet_names:  # see all sheet names
            df=xls.parse(sheet)  # read a specific sheet to DataFrame
            text+=str(df.to_dict('records'))
        return text



    @staticmethod
    def read_html(file_path):
        text=""
        some_html_string=read_txt(file_path)
        text = ' '.join(BeautifulSoup(some_html_string, "html.parser").stripped_strings)
        return text

    @staticmethod
    def read_pptx(file_path):
        text=""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):  # Verifica se la forma contiene testo
                    text += shape.text + "\n"  # Aggiungi il testo con un ritorno a capo
        return text
     

    @staticmethod
    def split_text(text, chunk_size=500, overlap=50):
        chunks = []
        start = 0
        end = 0

        while end < len(text):
            end = start + chunk_size
            if end > len(text):
                end = len(text)
            chunks.append(text[start:end])
            start = end - overlap  # Overlap chunks

        return chunks


